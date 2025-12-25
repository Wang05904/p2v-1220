from ppt_parser import pptx_to_images
import os
import json
from pptx import Presentation
import win32com.client
import pythoncom
from config import IMG_DIR, TEMP_DIR  # 假设你已经有了config配置

import os
import json
from pptx import Presentation
from config import TEMP_DIR, IMG_DIR

def run_deletion_test(json_file_path, ppt_file_path):
    """
    根据 JSON 里的原始 XML ID，从 PPT 中物理删除元素
    """
    os.makedirs(TEMP_DIR, exist_ok=True)
    output_pptx = os.path.join(TEMP_DIR, "temp_ppt.pptx")

    # 1. 加载数据
    if not os.path.exists(json_file_path):
        print("❌ 找不到 JSON 文件")
        return
    
    with open(json_file_path, 'r', encoding='utf-8') as f:
        json_data = json.load(f)

    prs = Presentation(ppt_file_path)
    print(f"✅ 成功加载 PPT: {ppt_file_path}")

    # 2. 执行基于 XML 的精准删除
    for slide_data in json_data["slides"]:
        slide_num = int(slide_data["slide_number"])
        if slide_num > len(prs.slides):
            continue
            
        slide = prs.slides[slide_num - 1]
        # 获取 JSON 中定义的该页所有待删 ID
        target_ids = [str(el["id"]) for el in slide_data["animated_elements"]]
        
        # 使用 XPath 搜索所有 cNvPr 节点
        # 这是最稳健的方法，能穿透组合和层级
        for cnvpr in slide._element.xpath('.//p:cNvPr'):
            if cnvpr.get('id') in target_ids:
                # 查找到对应的父级容器（即图片或形状节点）并移除
                shape_elm = cnvpr.getparent().getparent()
                if shape_elm is not None:
                    shape_elm.getparent().remove(shape_elm)
                    print(f"   - 第{slide_num}页: 已通过 ID {cnvpr.get('id')} 物理删除图片")

    # 3. 保存
    prs.save(output_pptx)
    print("-" * 50)
    print(f"🚀 任务完成！清理后的 PPT 已存至: {output_pptx}")
    pptx_to_images(output_pptx)

if __name__ == "__main__":
    run_deletion_test("extract_pic.json", "test.pptx")